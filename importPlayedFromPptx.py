import os
import shutil  # Import shutil for moving files
from pptx import Presentation
import re
import csv

def extract_text(pptx_file, csv_writer):
    # Load the Presentation from the .pptx file
    ppt = Presentation(pptx_file)

    # Extract text from the first and third slides (index 0 and 2)
    for index in [0, 2]:
        slide = ppt.slides[index]
        for shape in slide.shapes:
            if shape.has_text_frame and len(shape.text_frame.paragraphs) >= 5:
                for paragraph in shape.text_frame.paragraphs:
                    # Join all the runs in the paragraph
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    # Split the text into artist, title, and release year
                    parts = re.split(' – |- ', paragraph_text)
                    if len(parts) == 2:
                        artist, rest = parts
                        # Use a regular expression to extract the release year
                        match = re.search(r'\(?(\d{4})\)?$', rest)
                        if match:
                            release_year = match.group(1)
                            title = rest[:match.start()].strip()
                        else:
                            title = rest
                            release_year = ''
                        # Write the data to the CSV file
                        csv_writer.writerow([artist, title, release_year])
                        # Print the extracted data to console
                        print(f"Artist: {artist}, Title: {title}, Release Year: {release_year}")
                        
# Directory containing the .pptx files
directory = 'import/pptx/'
# Get the current working directory
current_directory = os.getcwd()
# Define the imported directory path
imported_directory = os.path.join(current_directory, 'import/pptx/imported')  

# Open the output file
with open(os.path.join(current_directory, 'usedSongs.tsv'), 'a', newline='') as output_file:
    tsv_writer = csv.writer(output_file, delimiter='\t')
    # Write the header row only if the file is empty - cur 172 as of game 18
    if output_file.tell() == 0:
        tsv_writer.writerow(['Artist', 'Title', 'Release Year'])

    # Get a list of all .pptx files in the directory and sort them    
    pptx_files = sorted([f for f in os.listdir(directory) if f.endswith('.pptx')])

    # Extract text from each .pptx file
    for pptx_file in pptx_files:
        extract_text(os.path.join(directory, pptx_file), tsv_writer)
        # Move the processed .pptx file to the /imported directory
        shutil.move(os.path.join(directory, pptx_file), os.path.join(imported_directory, pptx_file))